import pptx

def parse_ppt_content(file_path):
    ppt = pptx.Presentation(file_path)
    slides_text = {}
    for index, slide in enumerate(ppt.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)
            elif shape.has_table:
                for cell in shape.table.iter_cells():
                    slide_text.append(cell.text)
            elif shape.shape_type == 14 and not isinstance(shape, pptx.shapes.placeholder.PlaceholderPicture):  # 检查是否为文本框
                txBody = shape.text_frame
                for paragraph in txBody.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)
            elif shape.shape_type == 6:  # 检查是否为组合
                def recursive_get_text(subshape):
                    if subshape.has_text_frame:
                        for paragraph in subshape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_text.append(run.text)
                    elif subshape.shape_type == 6:  # 如果子形状也是组合
                        for subsubshape in subshape.shapes:
                            recursive_get_text(subsubshape)  # 递归获取文字

                for subshape in shape.shapes:
                    recursive_get_text(subshape)
        slides_text[index] = slide_text
    return slides_text

r = parse_ppt_content(r"")
for key in r:
    print(key, r[key])
